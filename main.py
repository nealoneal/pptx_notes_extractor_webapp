from flask import (Flask, render_template, send_file,
                   request, flash, redirect, url_for, send_from_directory)
import os
from werkzeug.utils import secure_filename
from pptx import Presentation
from docx import Document
from io import BytesIO, StringIO


UPLOAD_FOLDER = 'uploads'
ALLOWED_EXTENSIONS = ['pptx']

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['SECRET_KEY'] = "cccccclcdgvkdtluuvelgjcehtguhnickhdnrckneehf"


def extract_speaker_notes(original_file):
    global word_doc
    ppt = Presentation(original_file)
    notes = ""
    filename = original_file.name
    for page, slide in enumerate(ppt.slides):
        note = slide.notes_slide.notes_text_frame.text
        notes += (f"<Slide {int(page)+1}>"+"\n")
        notes += (note+"\n")
    print(type(notes))
    print(notes)
    word_doc = Document()
    word_doc.add_heading(f'Speaker Notes - {filename.split("/")[-1]}', 0)
    p = word_doc.add_paragraph(notes)
    files = [('Word Document', '*.docx')]
    word_doc.save(filename+".docx")


@app.route('/uploads/')
def download_file():
    return send_file(path_or_file=word_doc, as_attachment=True, download_name="speaker_notes.docx",
                     mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation")


def allowed_file(filename):
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS


@app.route('/', methods=['GET', 'POST'])
def upload_file():
    if request.method == 'POST':
        # check if the post request has the file part
        if 'file' not in request.files:
            print('No file part')

            return redirect(request.url)
        file = request.files['file']
        # If the user does not select a file, the browser submits an
        # empty file without a filename.
        if file.filename == '':
            flash('No selected file')
            return redirect(request.url)
        if file and allowed_file(file.filename):
            filename = secure_filename(file.filename)
            ppt = Presentation(file)
            notes = ""
            filename = file.filename
            print(filename)
            for page, slide in enumerate(ppt.slides):
                note = slide.notes_slide.notes_text_frame.text
                notes += (f"<Slide {int(page) + 1}>" + "\n")
                notes += (note + "\n")
            # print(type(notes))
            # print(notes)
            word_doc = Document()
            word_doc.add_heading(f'Speaker Notes - {filename.split(".")[0]}',0)
            p = word_doc.add_paragraph(notes)
            files = [('Word Document', '*.docx')]
            file_to_send = BytesIO()
            word_doc.save(file_to_send)
            file_to_send.seek(0)
            # file.save(os.path.join(app.config['UPLOAD_FOLDER'], filename))
            return send_file(path_or_file=file_to_send, as_attachment=True,
                             download_name=f"{filename.split('.')[0]}_speaker_notes.docx",
                            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation")

    return render_template('index.html')




# @app.route('/', methods=['GET', 'POST'])
# def upload_file():
#     if request.method == 'POST':
#         file = request.files['file']
#         if file and allowed_file(file.filename):
#             filename = secure_filename(file.filename)
#             file.save(os.path.join(app.config['UPLOAD_FOLDER'], filename))
#
#             ## snippet to read code below
#             file.stream.seek(0) # seek to the beginning of file
#             myfile = file.file # will point to tempfile itself
#             dataframe = pd.read_csv(myfile)
#             ## end snippet
#
#             return "yatta"
#         else:
#             return "file not allowed"
#
#     return render_template("index.html")
#

if __name__ == '__main__':
    app.run(port=5080, debug=True)
